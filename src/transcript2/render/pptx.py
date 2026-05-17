"""python-pptx renderer — clean Japanese consulting aesthetic.

16:9 deck, dark navy accent, generous whitespace, Japanese-safe font.
Frames (real or generated) are placed on the right column when present.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from ..schema import Deck, Slide

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
SLATE = RGBColor(0x1C, 0x3D, 0x5A)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x55, 0x55)

JP_FONT = "Yu Gothic"

SW, SH = Inches(13.333), Inches(7.5)


def _set_font(run, *, size, bold=False, color=NAVY):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = JP_FONT


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def _accent_bar(slide):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.18), SH)
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()


def _title_slide(prs, deck: Deck):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    s.shapes.add_shape(1, 0, Inches(5.4), SW, Inches(0.06)).fill.solid()
    tf = _box(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.2))
    p = tf.paragraphs[0]
    _set_font(p.add_run(), size=40, bold=True, color=WHITE)
    p.runs[0].text = deck.title or "プレゼンテーション"
    if deck.subtitle:
        p2 = tf.add_paragraph()
        _set_font(p2.add_run(), size=20, color=LIGHT)
        p2.runs[0].text = deck.subtitle
    foot = _box(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5))
    fp = foot.paragraphs[0]
    src = deck.meta.url if deck.meta else ""
    _set_font(fp.add_run(), size=11, color=GREY)
    fp.runs[0].text = f"出典: {src}"


def _content_slide(prs, slide: Slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _accent_bar(s)

    # Implication-driven headline.
    htf = _box(s, Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.2))
    hp = htf.paragraphs[0]
    _set_font(hp.add_run(), size=26, bold=True, color=NAVY)
    hp.runs[0].text = slide.title

    # Key message ribbon.
    if slide.message:
        mtf = _box(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.7))
        mp = mtf.paragraphs[0]
        _set_font(mp.add_run(), size=15, bold=True, color=GOLD)
        mp.runs[0].text = slide.message

    has_img = bool(slide.image_path and Path(slide.image_path).exists())
    body_w = Inches(7.4) if has_img else Inches(12.1)

    btf = _box(s, Inches(0.6), Inches(2.5), body_w, Inches(4.4))
    btf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for b in slide.bullets:
        para = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        _set_font(para.add_run(), size=17, color=SLATE)
        para.runs[0].text = f"●  {b}"
        para.space_after = Pt(10)

    if has_img:
        try:
            s.shapes.add_picture(
                slide.image_path, Inches(8.3), Inches(2.5),
                width=Inches(4.4), height=Inches(3.3),
            )
        except Exception:
            pass


def _section_slide(prs, slide: Slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, SLATE)
    tf = _box(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.6))
    p = tf.paragraphs[0]
    _set_font(p.add_run(), size=32, bold=True, color=WHITE)
    p.runs[0].text = slide.title
    if slide.message:
        p2 = tf.add_paragraph()
        _set_font(p2.add_run(), size=16, color=LIGHT)
        p2.runs[0].text = slide.message


def render_pptx(deck: Deck, out_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    for i, slide in enumerate(deck.slides):
        if slide.layout == "title" or (i == 0 and slide.layout != "section"):
            if i == 0:
                _title_slide(prs, deck)
            else:
                _section_slide(prs, slide)
        elif slide.layout in ("section", "closing"):
            _section_slide(prs, slide)
        else:
            _content_slide(prs, slide)

        # Speaker notes.
        if slide.speaker_notes:
            notes = prs.slides[-1].notes_slide.notes_text_frame
            notes.text = slide.speaker_notes

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
